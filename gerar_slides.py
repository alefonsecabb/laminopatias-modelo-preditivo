"""
Gera Apresentacao_Laminopatias_ML.pptx no diretório do projeto.
Requer: pip install python-pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm

BASE_DIR = os.path.dirname(os.path.abspath(__file__))

# Cores
AZUL_ESCURO  = RGBColor(0x00, 0x33, 0x66)
AZUL_MEDIO   = RGBColor(0x00, 0x5B, 0x9A)
AZUL_CLARO   = RGBColor(0xD6, 0xE4, 0xF0)
BRANCO       = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_TEXTO  = RGBColor(0x33, 0x33, 0x33)
VERMELHO     = RGBColor(0xC0, 0x00, 0x00)
VERDE        = RGBColor(0x00, 0x70, 0x00)


def nova_apresentacao():
    prs = Presentation()
    prs.slide_width  = Cm(33.87)
    prs.slide_height = Cm(19.05)
    return prs


def fundo_azul(slide, prs):
    from pptx.util import Emu
    fundo = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, prs.slide_width, prs.slide_height
    )
    fundo.fill.solid()
    fundo.fill.fore_color.rgb = AZUL_ESCURO
    fundo.line.fill.background()
    fundo.zorder = 0


def fundo_branco_com_barra(slide, prs):
    """Fundo branco + barra azul no topo (2 cm)."""
    barra = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Cm(2.2))
    barra.fill.solid()
    barra.fill.fore_color.rgb = AZUL_ESCURO
    barra.line.fill.background()


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=CINZA_TEXTO,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def add_bullet_slide(slide, prs, titulo, bullets, subtitulo=None):
    fundo_branco_com_barra(slide, prs)
    # Título na barra
    add_textbox(slide, titulo,
                Cm(0.5), Cm(0.2), Cm(32), Cm(1.8),
                font_size=22, bold=True, color=BRANCO)
    y = Cm(2.8)
    if subtitulo:
        add_textbox(slide, subtitulo,
                    Cm(1), y, Cm(31), Cm(0.8),
                    font_size=14, bold=False, color=AZUL_MEDIO, italic=True)
        y += Cm(1.0)
    for bullet in bullets:
        nivel  = bullet.get("nivel", 0)
        texto  = bullet["texto"]
        cor    = bullet.get("cor", CINZA_TEXTO)
        negrito = bullet.get("bold", False)
        indent = Cm(1 + nivel * 0.8)
        add_textbox(slide, ("• " if nivel == 0 else "   ◦ ") + texto,
                    Cm(1) + indent, y, Cm(30) - indent, Cm(0.75),
                    font_size=16, bold=negrito, color=cor)
        y += Cm(0.72)


def add_imagem(slide, nome_arquivo, left, top, width=None, height=None):
    caminho = os.path.join(BASE_DIR, nome_arquivo)
    if not os.path.exists(caminho):
        print(f"[AVISO] Figura não encontrada: {caminho}")
        return
    if width and height:
        slide.shapes.add_picture(caminho, left, top, width, height)
    elif width:
        slide.shapes.add_picture(caminho, left, top, width=width)
    elif height:
        slide.shapes.add_picture(caminho, left, top, height=height)
    else:
        slide.shapes.add_picture(caminho, left, top)


def add_tabela(slide, dados, left, top, width, height,
               header_bg=AZUL_ESCURO, header_fg=BRANCO,
               row_alt=AZUL_CLARO):
    from pptx.util import Pt
    rows = len(dados)
    cols = len(dados[0])
    tabela = slide.shapes.add_table(rows, cols, left, top, width, height).table
    col_w = width // cols
    for c in range(cols):
        tabela.columns[c].width = col_w
    for r, linha in enumerate(dados):
        for c, cell_text in enumerate(linha):
            cell = tabela.cell(r, c)
            cell.text = str(cell_text)
            tf = cell.text_frame
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            tf.paragraphs[0].runs[0].font.size = Pt(13)
            tf.paragraphs[0].runs[0].font.bold = (r == 0)
            tf.paragraphs[0].runs[0].font.color.rgb = header_fg if r == 0 else CINZA_TEXTO
            fill = cell.fill
            fill.solid()
            if r == 0:
                fill.fore_color.rgb = header_bg
            elif r % 2 == 0:
                fill.fore_color.rgb = row_alt
            else:
                fill.fore_color.rgb = BRANCO


# ─────────────────────────────────────────────
# SLIDES
# ─────────────────────────────────────────────

def slide_capa(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fundo_azul(slide, prs)

    add_textbox(slide,
                "Expansão do Modelo Preditivo para\nTaquiarritmias Ventriculares em Laminopatias",
                Cm(2), Cm(3.5), Cm(30), Cm(4),
                font_size=30, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)

    add_textbox(slide,
                "Proposta de Inclusão de Fibrose Miocárdica (RMC) e Duração do QRS",
                Cm(2), Cm(7.5), Cm(30), Cm(1.2),
                font_size=18, bold=False, color=AZUL_CLARO, align=PP_ALIGN.CENTER, italic=True)

    add_textbox(slide,
                "Alexandre da Fonseca  ·  Paloma de Almeida Taboada",
                Cm(2), Cm(9.5), Cm(30), Cm(0.9),
                font_size=16, bold=False, color=BRANCO, align=PP_ALIGN.CENTER)

    add_textbox(slide,
                "Disciplina: Aprendizado de Máquina  |  FATEC Jundiaí  |  2026",
                Cm(2), Cm(10.5), Cm(30), Cm(0.8),
                font_size=13, color=AZUL_CLARO, align=PP_ALIGN.CENTER)

    add_textbox(slide,
                "Orientador: Prof. Argemiro Pentian Junior",
                Cm(2), Cm(11.5), Cm(30), Cm(0.7),
                font_size=12, color=AZUL_CLARO, align=PP_ALIGN.CENTER, italic=True)


def slide_laminopatias(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullet_slide(slide, prs,
        "O que são Laminopatias?",
        [
            {"texto": "Doenças genéticas causadas por mutações no gene LMNA", "bold": True},
            {"texto": "Gene LMNA codifica as proteínas Lamina A/C — estrutura do núcleo celular", "nivel": 1},
            {"texto": "Afetam principalmente o coração e o músculo esquelético", "nivel": 1},
            {"texto": "Principal risco cardíaco: Taquiarritmia Ventricular Ameaçadora à Vida (LTVTA)", "bold": True},
            {"texto": "LTVTA = TV sustentada, FV ou morte súbita cardíaca", "nivel": 1},
            {"texto": "Incidência acumulada em 5 anos: ~18% nos portadores de mutação LMNA", "nivel": 1},
            {"texto": "Prevenção: cardiodesfibrilador implantável (CDI) — mas indicação é difícil de definir", "bold": True},
            {"texto": "Necessidade de estratificação de risco precisa para decidir quem implanta o CDI", "nivel": 1},
        ],
        subtitulo="Contexto Clínico"
    )


def slide_wahbi(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fundo_branco_com_barra(slide, prs)
    add_textbox(slide, "Modelo de Wahbi et al. (2019) — Referência Atual",
                Cm(0.5), Cm(0.2), Cm(32), Cm(1.8),
                font_size=22, bold=True, color=BRANCO)

    add_textbox(slide, "Fine-Gray (risco sub-distribuição) para riscos competidores — LTVTA vs. Óbito",
                Cm(1), Cm(2.5), Cm(31), Cm(0.8),
                font_size=14, color=AZUL_MEDIO, italic=True)

    dados = [
        ["Variável", "Tipo", "Hazard Ratio", "Interpretação"],
        ["Sexo masculino",       "Binária",    "2,43",         "↑ 143% risco"],
        ["Mutação não-missense", "Binária",    "2,54",         "↑ 154% risco"],
        ["Bloqueio AV ≥1º grau","Binária",    "2,87",         "↑ 187% risco"],
        ["TVNS",                 "Binária",    "2,75",         "↑ 175% risco"],
        ["FEVE (por 10% ↑)",    "Contínua",  "0,64",         "↓ 36% risco"],
    ]
    add_tabela(slide, dados,
               Cm(1), Cm(3.5), Cm(31), Cm(6.5))

    add_textbox(slide,
                "C-index: 0,776 (derivação)  ·  0,800 (validação externa)",
                Cm(1), Cm(10.5), Cm(31), Cm(0.8),
                font_size=15, bold=True, color=AZUL_ESCURO, align=PP_ALIGN.CENTER)

    add_textbox(slide,
                "Coorte: 311 pacientes europeus com mutação LMNA confirmada  ·  5 anos de seguimento",
                Cm(1), Cm(11.4), Cm(31), Cm(0.7),
                font_size=12, color=CINZA_TEXTO, align=PP_ALIGN.CENTER, italic=True)


def slide_proposta(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullet_slide(slide, prs,
        "Nossa Proposta — Duas Novas Variáveis",
        [
            {"texto": "LGE por Ressonância Magnética Cardíaca (RMC)",
             "bold": True, "cor": AZUL_ESCURO},
            {"texto": "LGE = Late Gadolinium Enhancement → detecta fibrose miocárdica", "nivel": 1},
            {"texto": "Fibrose = substrato para reentrada e arritmias ventriculares", "nivel": 1},
            {"texto": "Hipótese: HR = 2,0 – 4,0", "nivel": 1, "cor": AZUL_MEDIO},
            {"texto": "Duração do QRS (ECG de superfície)",
             "bold": True, "cor": AZUL_ESCURO},
            {"texto": "QRS ≥120 ms indica distúrbio de condução intraventricular", "nivel": 1},
            {"texto": "Associado a disfunção elétrica e mecânica do VE", "nivel": 1},
            {"texto": "Hipótese: HR = 1,10 – 1,30 a cada 10 ms", "nivel": 1, "cor": AZUL_MEDIO},
            {"texto": "Objetivo: verificar se as novas variáveis melhoram o poder discriminativo do modelo",
             "bold": True, "cor": VERDE},
        ],
        subtitulo="Hipóteses de Pesquisa"
    )


def slide_dataset(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullet_slide(slide, prs,
        "Metodologia — Dataset Sintético Parametrizado",
        [
            {"texto": "n = 600 pacientes simulados  (numpy.random.seed = 42)", "bold": True},
            {"texto": "Parâmetros calibrados a partir da coorte Wahbi et al. (2019)", "nivel": 1},
            {"texto": "Geração de eventos: distribuição Weibull com riscos competidores"},
            {"texto": "Evento 1 — LTVTA: 12,8% (77 pacientes)  →  ~18% em 5 anos", "nivel": 1, "cor": VERMELHO},
            {"texto": "Evento 2 — Óbito competidor: 9,5%", "nivel": 1},
            {"texto": "Censura administrativa: 77,7%  ·  5% perda aleatória de seguimento", "nivel": 1},
            {"texto": "Calibração automática da taxa base (scipy.optimize.brentq) para atingir 18% em 5 anos"},
            {"texto": "Coeficientes de simulação derivados de ln(HR) da literatura", "nivel": 1},
            {"texto": "Tempo médio de seguimento: 4,33 ± 1,3 anos"},
        ],
        subtitulo="Dados Simulados — validados contra literatura"
    )


def slide_eda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fundo_branco_com_barra(slide, prs)
    add_textbox(slide, "Análise Exploratória — Prevalências vs. Literatura",
                Cm(0.5), Cm(0.2), Cm(32), Cm(1.8),
                font_size=22, bold=True, color=BRANCO)

    add_imagem(slide, "fig2_prevalencias.png",
               Cm(1), Cm(2.3), width=Cm(31), height=Cm(15.5))


def slide_modelo_ml(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullet_slide(slide, prs,
        "Modelo de ML — Cox com Riscos Competidores",
        [
            {"texto": "Abordagem: Regressão de Cox causa-específica", "bold": True},
            {"texto": "Biblioteca: scikit-survival (Python)  ·  CoxPHSurvivalAnalysis", "nivel": 1},
            {"texto": "Modelagem separada para cada causa de evento (LTVTA / Óbito)", "nivel": 1},
            {"texto": "Por que não Fine-Gray exato?"},
            {"texto": "Fine-Gray exato requer pacote cmprsk do R", "nivel": 1},
            {"texto": "Cox causa-específica é metodologicamente válido para comparação relativa de C-index", "nivel": 1},
            {"texto": "Mesma população → comparação justa entre Modelo Original e Expandido", "nivel": 1},
            {"texto": "Avaliação de desempenho"},
            {"texto": "C-index de Harrell (discriminação): capacidade de ordenar pacientes por risco", "nivel": 1},
            {"texto": "NRI (Net Reclassification Index): melhora de classificação com limiar de 7%", "nivel": 1},
            {"texto": "Análise de sensibilidade: variar HR_LGE ∈ [2,0;4,0] e HR_QRS ∈ [1,10;1,30]", "nivel": 1},
        ],
    )


def slide_finegray_vs_cox(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fundo_branco_com_barra(slide, prs)
    add_textbox(slide, "Fine-Gray vs. Cox Causa-Específica",
                Cm(0.5), Cm(0.2), Cm(32), Cm(1.8),
                font_size=22, bold=True, color=BRANCO)

    # ── Coluna esquerda: Cox ───────────────────────────────────────────
    add_textbox(slide, "Cox Causa-Específica",
                Cm(0.8), Cm(2.4), Cm(15), Cm(0.7),
                font_size=16, bold=True, color=AZUL_ESCURO)

    linhas_cox = [
        "Modela a taxa condicional de LTVTA",
        "dado que o paciente ainda está vivo",
        "",
        "• Competing events → censurados",
        "  (saem do conjunto de risco)",
        "",
        "• Pergunta respondida:",
        '  "Quais fatores aumentam a TAXA',
        '   de LTVTA entre os ainda vivos?"',
        "",
        "• Melhor para inferência etiológica",
        "• Disponível no scikit-survival ✓",
    ]
    y = Cm(3.2)
    for l in linhas_cox:
        add_textbox(slide, l, Cm(0.8), y, Cm(15.5), Cm(0.55),
                    font_size=13, color=CINZA_TEXTO)
        y += Cm(0.52)

    # ── Divisor vertical ──────────────────────────────────────────────
    div = slide.shapes.add_shape(1, Cm(17), Cm(2.4), Cm(0.05), Cm(11))
    div.fill.solid()
    div.fill.fore_color.rgb = AZUL_CLARO
    div.line.fill.background()

    # ── Coluna direita: Fine-Gray ──────────────────────────────────────
    add_textbox(slide, "Fine-Gray (Subdistribution Hazard)",
                Cm(17.5), Cm(2.4), Cm(16), Cm(0.7),
                font_size=16, bold=True, color=AZUL_ESCURO)

    linhas_fg = [
        "Modela diretamente a probabilidade",
        "acumulada real de LTVTA — CIF F(t)",
        "",
        "• Competing events → ficam no risco",
        '  com peso → 0  ("pacientes mortos-vivos")',
        "",
        "• Pergunta respondida:",
        '  "Qual a CHANCE REAL de LTVTA',
        '   em 5 anos, considerando morte?"',
        "",
        "• Melhor para predição absoluta de risco",
        "• Requer cmprsk (R) — não disponível em Python",
    ]
    y = Cm(3.2)
    for l in linhas_fg:
        add_textbox(slide, l, Cm(17.5), y, Cm(15.5), Cm(0.55),
                    font_size=13, color=CINZA_TEXTO)
        y += Cm(0.52)

    # ── Rodapé: por que Cox é válido aqui ─────────────────────────────
    rodape = slide.shapes.add_shape(1, Cm(0.5), Cm(14), Cm(32.5), Cm(4.5))
    rodape.fill.solid()
    rodape.fill.fore_color.rgb = AZUL_CLARO
    rodape.line.fill.background()

    add_textbox(slide,
                "Por que Cox causa-específica é válido neste estudo?",
                Cm(1), Cm(14.2), Cm(31), Cm(0.7),
                font_size=14, bold=True, color=AZUL_ESCURO)
    add_textbox(slide,
                "Objetivo = COMPARAR dois modelos na mesma população, não estimar risco absoluto.\n"
                "O C-index mede apenas a capacidade de ORDENAR pacientes por risco — e essa ordenação\n"
                "é equivalente entre Cox e Fine-Gray na mesma coorte. A escolha não altera a conclusão.",
                Cm(1), Cm(15.0), Cm(31.5), Cm(3.2),
                font_size=13, color=CINZA_TEXTO)


def slide_hr(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fundo_branco_com_barra(slide, prs)
    add_textbox(slide, "Resultados — Hazard Ratios Estimados",
                Cm(0.5), Cm(0.2), Cm(32), Cm(1.8),
                font_size=22, bold=True, color=BRANCO)

    dados = [
        ["Variável",                "HR Wahbi (ref.)", "HR Modelo Orig.", "HR Modelo Exp."],
        ["Sexo masculino",           "2,43",            "2,41",            "2,75"],
        ["Mutação não-missense",     "2,54",            "2,47",            "2,48"],
        ["Bloqueio AV",              "2,87",            "3,32",            "2,59"],
        ["TVNS",                     "2,75",            "2,60",            "2,63"],
        ["FEVE (por 10% ↑)",        "0,64",            "0,97",            "0,97"],
        ["LGE presente [NOVO]",      "—",               "—",               "2,26 ✓"],
        ["QRS / 10 ms [NOVO]",       "—",               "—",               "1,23 ✓"],
    ]
    add_tabela(slide, dados,
               Cm(1), Cm(2.5), Cm(22), Cm(8.5))

    add_imagem(slide, "fig6_tornado.png",
               Cm(23.5), Cm(2.5), width=Cm(10), height=Cm(10.5))

    add_textbox(slide,
                "✓ Ambas as hipóteses confirmadas nas faixas esperadas",
                Cm(1), Cm(11.3), Cm(22), Cm(0.8),
                font_size=14, bold=True, color=VERDE)


def slide_desempenho(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fundo_branco_com_barra(slide, prs)
    add_textbox(slide, "Resultados — Desempenho Discriminativo (C-index)",
                Cm(0.5), Cm(0.2), Cm(32), Cm(1.8),
                font_size=22, bold=True, color=BRANCO)

    add_textbox(slide,
                "Modelo Original (5 var.)     C-index = 0,768",
                Cm(1), Cm(2.5), Cm(15), Cm(0.8),
                font_size=16, bold=True, color=CINZA_TEXTO)

    add_textbox(slide,
                "Modelo Expandido (7 var.)   C-index = 0,794   (+0,026)",
                Cm(1), Cm(3.3), Cm(20), Cm(0.8),
                font_size=16, bold=True, color=VERDE)

    add_textbox(slide,
                "NRI Total: +2,6%  (limiar 7% de risco em 5 anos)",
                Cm(1), Cm(4.1), Cm(22), Cm(0.7),
                font_size=14, color=CINZA_TEXTO)

    add_imagem(slide, "fig4_desempenho.png",
               Cm(1), Cm(5.0), width=Cm(32), height=Cm(12.5))


def slide_sensibilidade(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fundo_branco_com_barra(slide, prs)
    add_textbox(slide, "Análise de Sensibilidade — Robustez do Resultado",
                Cm(0.5), Cm(0.2), Cm(32), Cm(1.8),
                font_size=22, bold=True, color=BRANCO)

    add_textbox(slide,
                "C-index do Modelo Expandido permanece acima do Original em TODA a faixa de HR testada",
                Cm(1), Cm(2.5), Cm(32), Cm(0.8),
                font_size=14, bold=True, color=AZUL_MEDIO)

    add_imagem(slide, "fig5_sensibilidade.png",
               Cm(1), Cm(3.4), width=Cm(32), height=Cm(13.5))


def slide_limitacoes(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullet_slide(slide, prs,
        "Limitações do Estudo",
        [
            {"texto": "Dados sintéticos — não substituem coorte clínica real", "bold": True, "cor": VERMELHO},
            {"texto": "Gerados a partir de parâmetros da literatura, não pacientes reais", "nivel": 1},
            {"texto": "n = 77 eventos (LTVTA) — poder amostral limitado", "bold": True, "cor": VERMELHO},
            {"texto": "Hipótese de C-index > 0,82 não confirmada (obtido: 0,794)", "nivel": 1},
            {"texto": "Hipótese de NRI ≥ 10% não confirmada (obtido: +2,6%)", "nivel": 1},
            {"texto": "Implementação Cox causa-específica vs. Fine-Gray exato"},
            {"texto": "Fine-Gray exato exigiria o pacote cmprsk do R", "nivel": 1},
            {"texto": "Diferença metodológica não invalida comparação relativa entre modelos", "nivel": 1},
            {"texto": "Próximos passos para validação real"},
            {"texto": "Aplicar a coorte clínica prospectiva com dados de RMC e ECG", "nivel": 1},
            {"texto": "Validação externa independente", "nivel": 1},
        ],
    )


def slide_conclusoes(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullet_slide(slide, prs,
        "Conclusões e Implicações Clínicas",
        [
            {"texto": "LGE (fibrose por RMC) e QRS se confirmam como preditores relevantes",
             "bold": True, "cor": VERDE},
            {"texto": "HR LGE = 2,26 ✓  (hipótese: 2,0–4,0)", "nivel": 1},
            {"texto": "HR QRS = 1,23/10ms ✓  (hipótese: 1,10–1,30)", "nivel": 1},
            {"texto": "Melhora consistente do C-index: 0,768 → 0,794  (+3,4%)",
             "bold": True, "cor": VERDE},
            {"texto": "Confirmado em 100% dos cenários da análise de sensibilidade", "nivel": 1},
            {"texto": "Implicação clínica direta"},
            {"texto": "RMC deveria ser incluída nos protocolos de estratificação de LTVTA em laminopatias",
             "nivel": 1, "bold": True},
            {"texto": "QRS ≥120 ms é marcador de risco incremental de baixo custo", "nivel": 1},
            {"texto": "Análise computacional disponível no Google Colab (link no artigo)", "bold": False},
            {"texto": "Reprodutível: seed=42, sem dados de pacientes reais", "nivel": 1},
        ],
        subtitulo="Obrigado!"
    )


def main():
    prs = nova_apresentacao()
    slide_capa(prs)
    slide_laminopatias(prs)
    slide_wahbi(prs)
    slide_proposta(prs)
    slide_dataset(prs)
    slide_eda(prs)
    slide_modelo_ml(prs)
    slide_finegray_vs_cox(prs)
    slide_hr(prs)
    slide_desempenho(prs)
    slide_sensibilidade(prs)
    slide_limitacoes(prs)
    slide_conclusoes(prs)

    out = os.path.join(BASE_DIR, "Apresentacao_Laminopatias_ML.pptx")
    prs.save(out)
    print(f"Apresentação gerada: {out}")
    print(f"Total de slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
